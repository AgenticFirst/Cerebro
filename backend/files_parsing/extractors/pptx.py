"""Extract markdown from .pptx via python-pptx. One H2 per slide."""

from __future__ import annotations

from pptx import Presentation  # type: ignore[import-untyped]
from pptx import __version__ as PPTX_VERSION  # type: ignore[import-untyped]


PARSER_NAME = "python-pptx"
PARSER_VERSION = PPTX_VERSION


def extract(file_path: str) -> str:
    pres = Presentation(file_path)
    parts: list[str] = []
    for idx, slide in enumerate(pres.slides, start=1):
        block = [f"## Slide {idx}"]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    block.append(text)
        if len(block) > 1:
            parts.append("\n\n".join(block))
    return ("\n\n".join(parts) or "").strip() + "\n"
